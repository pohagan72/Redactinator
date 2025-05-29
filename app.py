#foo 
import os
import io
import tempfile
from flask import Flask, request, render_template, send_file, flash, redirect, url_for
from werkzeug.utils import secure_filename
from docx import Document
from docx.shared import RGBColor
from docx.enum.text import WD_COLOR_INDEX # For highlighting
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor # Need specific import
from pptx.util import Pt

from presidio_analyzer import AnalyzerEngine
from presidio_analyzer.nlp_engine import NlpEngineProvider

# --- Presidio Configuration ---
# Set up NLP engine provider (using Spacy)
provider = NlpEngineProvider(nlp_configuration={
    "nlp_engine_name": "spacy",
    "models": [{"lang_code": "en", "model_name": "en_core_web_lg"}] # Use the downloaded model
})
analyzer = AnalyzerEngine(nlp_engine=provider.create_engine(), supported_languages=["en"])
# -----------------------------

# --- Flask Configuration ---
app = Flask(__name__)
app.config['SECRET_KEY'] = os.urandom(24)  # Needed for flashing messages
ALLOWED_EXTENSIONS = {'docx', 'pptx'}
# --------------------------

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# --- Redaction Logic ---

def redact_word_document(file_stream):
    """
    Reads a .docx file stream, identifies PII using Presidio,
    and applies black highlighting to the identified text.
    Returns a stream containing the redacted document.
    """
    try:
        document = Document(file_stream)
        redacted_runs_count = 0

        # --- Process Paragraphs ---
        for para in document.paragraphs:
            if not para.text.strip(): # Skip empty paragraphs
                continue

            try:
                analyzer_results = analyzer.analyze(text=para.text, language='en')
            except Exception as e:
                print(f"Error analyzing paragraph text: {para.text[:50]}... Error: {e}")
                continue # Skip paragraph if analysis fails

            # Keep track of offset within the paragraph's runs
            current_offset = 0
            run_index = 0
            runs = para.runs

            # Create a list of (start, end) tuples for redaction ranges
            redaction_ranges = sorted([(res.start, res.end) for res in analyzer_results])
            range_index = 0

            while run_index < len(runs) and range_index < len(redaction_ranges):
                run = runs[run_index]
                run_len = len(run.text)
                run_start = current_offset
                run_end = current_offset + run_len

                redact_start, redact_end = redaction_ranges[range_index]

                # Check for overlap between current run and current redaction range
                overlap_start = max(run_start, redact_start)
                overlap_end = min(run_end, redact_end)

                if overlap_start < overlap_end:
                    # Apply redaction (black highlight) to the entire run if it overlaps
                    # More precise redaction (splitting runs) is complex with python-docx
                    run.font.highlight_color = WD_COLOR_INDEX.BLACK
                    # Optionally, also make text black (though highlight often suffices)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    redacted_runs_count += 1

                    # If the redaction range continues beyond this run, keep the same range_index
                    # If the run continues beyond the redaction range, we still move to the next run for simplicity
                    if redact_end <= run_end:
                        range_index += 1 # Move to the next redaction range

                # Move to the next run or stay if redaction spans multiple runs
                if run_end >= redact_end and overlap_start >= overlap_end:
                     # If current run ends after or at the redaction end, and no overlap was found for this run
                     # move to the next redaction range (if run_start was already > redact_end)
                     if run_start >= redact_end:
                         range_index += 1
                     else: # Standard case: move run pointer
                        current_offset += run_len
                        run_index += 1
                else: # Standard case: move run pointer
                    current_offset += run_len
                    run_index += 1


        # --- Process Tables ---
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    # Recursively process paragraphs within the cell
                     for para in cell.paragraphs:
                        if not para.text.strip(): continue
                        try:
                           analyzer_results = analyzer.analyze(text=para.text, language='en')
                        except Exception as e:
                           print(f"Error analyzing table cell text: {para.text[:50]}... Error: {e}")
                           continue

                        current_offset = 0
                        run_index = 0
                        runs = para.runs
                        redaction_ranges = sorted([(res.start, res.end) for res in analyzer_results])
                        range_index = 0

                        # (Duplicate the run processing logic from above for cell paragraphs)
                        while run_index < len(runs) and range_index < len(redaction_ranges):
                            run = runs[run_index]
                            run_len = len(run.text)
                            run_start = current_offset
                            run_end = current_offset + run_len
                            redact_start, redact_end = redaction_ranges[range_index]
                            overlap_start = max(run_start, redact_start)
                            overlap_end = min(run_end, redact_end)

                            if overlap_start < overlap_end:
                                run.font.highlight_color = WD_COLOR_INDEX.BLACK
                                run.font.color.rgb = RGBColor(0, 0, 0)
                                redacted_runs_count += 1
                                if redact_end <= run_end:
                                    range_index += 1
                            if run_end >= redact_end and overlap_start >= overlap_end:
                                if run_start >= redact_end:
                                    range_index += 1
                                else:
                                    current_offset += run_len
                                    run_index += 1
                            else:
                                current_offset += run_len
                                run_index += 1


        print(f"Attempted to redact {redacted_runs_count} runs in Word document.")

        # Save redacted document to a BytesIO stream
        output_stream = io.BytesIO()
        document.save(output_stream)
        output_stream.seek(0)
        return output_stream

    except Exception as e:
        print(f"Error processing Word document: {e}")
        raise # Re-raise the exception to be caught in the Flask route

def redact_powerpoint_document(file_stream):
    """
    Reads a .pptx file stream, identifies PII using Presidio,
    and replaces the text in runs containing PII with '█' characters.
    Returns a stream containing the redacted document.

    Note: Redaction in PPTX is less precise due to library limitations.
          This version replaces the entire text of any run found containing PII.
    """
    try:
        presentation = Presentation(file_stream)
        redacted_runs_count = 0

        for i, slide in enumerate(presentation.slides):
            print(f"Processing Slide {i+1}/{len(presentation.slides)}")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text_frame = shape.text_frame
                for para_idx, para in enumerate(text_frame.paragraphs):
                    if not para.text.strip():
                        continue

                    try:
                        analyzer_results = analyzer.analyze(text=para.text, language='en')
                    except Exception as e:
                        print(f"Error analyzing paragraph text (Slide {i+1}, Shape {shape.shape_id}, Para {para_idx}): {para.text[:50]}... Error: {e}")
                        continue

                    current_offset = 0
                    run_index = 0
                    runs = para.runs
                    redaction_ranges = sorted([(res.start, res.end) for res in analyzer_results])
                    range_index = 0

                    while run_index < len(runs) and range_index < len(redaction_ranges):
                        run = runs[run_index]
                        # Handle potential empty runs added during previous edits
                        if not hasattr(run, 'text'):
                             run_index += 1
                             continue
                        run_len = len(run.text)
                        run_start = current_offset
                        run_end = current_offset + run_len

                        # If run has no length (can happen), skip it
                        if run_len == 0:
                            run_index += 1
                            continue

                        redact_start, redact_end = redaction_ranges[range_index]

                        # Check for overlap
                        overlap_start = max(run_start, redact_start)
                        overlap_end = min(run_end, redact_end)

                        if overlap_start < overlap_end:
                            # --- Redaction: Replace run text with blocks ---
                            original_text = run.text
                            redaction_char = '█'
                            # Make font black as well for better visual redaction
                            run.font.color.rgb = PptxRGBColor(0, 0, 0)
                             # Set font size if needed, otherwise it might shrink
                            if run.font.size is None and para.font.size is not None:
                                run.font.size = para.font.size
                            elif run.font.size is None:
                                run.font.size = Pt(11) # Default fallback size

                            run.text = redaction_char * len(original_text)
                            print(f"  - Redacted run (Slide {i+1}, Shape {shape.shape_id}, Para {para_idx}, Run {run_index}): '{original_text[:20]}...' -> '{run.text[:20]}...'")
                            redacted_runs_count += 1
                            # --- End Redaction ---

                            # Move to next redaction range if it ends within or at the end of this run
                            if redact_end <= run_end:
                                range_index += 1

                        # Logic to advance pointers
                        if run_end >= redact_end and overlap_start >= overlap_end:
                             if run_start >= redact_end:
                                 range_index += 1
                             else:
                                current_offset += run_len
                                run_index += 1
                        else:
                            current_offset += run_len
                            run_index += 1


        print(f"Attempted to redact {redacted_runs_count} runs in PowerPoint document.")

        # Save redacted presentation to a BytesIO stream
        output_stream = io.BytesIO()
        presentation.save(output_stream)
        output_stream.seek(0)
        return output_stream

    except Exception as e:
        print(f"Error processing PowerPoint document: {e}")
        import traceback
        traceback.print_exc() # Print detailed traceback for debugging pptx issues
        raise # Re-raise the exception


# --- Flask Routes ---

@app.route('/', methods=['GET'])
def index():
    """Renders the upload page."""
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    """Handles file upload, processing, and download."""
    if 'file' not in request.files:
        flash('No file part', 'error')
        return redirect(url_for('index'))

    file = request.files['file']

    if file.filename == '':
        flash('No selected file', 'error')
        return redirect(url_for('index'))

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_ext = filename.rsplit('.', 1)[1].lower()
        
        # Read file into memory (BytesIO) to avoid saving potentially sensitive intermediate file
        file_stream = io.BytesIO(file.read())
        file_stream.seek(0) # Reset stream position

        output_stream = None
        redacted_filename = f"redacted_{filename}"
        mime_type = None

        try:
            if file_ext == 'docx':
                print(f"Processing Word file: {filename}")
                output_stream = redact_word_document(file_stream)
                mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif file_ext == 'pptx':
                print(f"Processing PowerPoint file: {filename}")
                output_stream = redact_powerpoint_document(file_stream)
                mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

            if output_stream:
                print(f"Redaction complete. Sending file: {redacted_filename}")
                return send_file(
                    output_stream,
                    as_attachment=True,
                    download_name=redacted_filename,
                    mimetype=mime_type
                )
            else:
                # This case should ideally not be reached if exceptions are handled
                flash('File processing failed internally.', 'error')
                return redirect(url_for('index'))

        except Exception as e:
            flash(f'An error occurred during processing: {str(e)}', 'error')
            print(f"Error during processing route: {e}") # Log the error server-side
            import traceback
            traceback.print_exc()
            return redirect(url_for('index'))
        finally:
             # Ensure the input stream is closed
             file_stream.close()
             # output_stream is handled by send_file or closed implicitly if an error occurs before send_file

    else:
        flash('Invalid file type. Please upload .docx or .pptx files.', 'error')
        return redirect(url_for('index'))

# --- Main Execution ---
if __name__ == '__main__':
    # Use waitress or gunicorn for production deployments
    app.run(debug=True, host='0.0.0.0', port=5000) # Use debug=False in production